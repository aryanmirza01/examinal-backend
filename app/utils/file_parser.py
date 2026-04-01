"""
Parse PDF, DOCX, and PPTX files into page‑level text dicts.
Returns: List[{"text": str, "page": int | None}]
"""

import logging
from typing import List, Dict, Any

logger = logging.getLogger(__name__)


def parse_pdf(path: str) -> List[Dict[str, Any]]:
    import pdfplumber

    pages: List[Dict[str, Any]] = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            # Also try extracting tables
            tables = page.extract_tables() or []
            for table in tables:
                for row in table:
                    if row:
                        text += "\n" + " | ".join(str(cell or "") for cell in row)
            pages.append({"text": text, "page": i + 1})
    logger.info("Parsed PDF %s: %d pages", path, len(pages))
    return pages


def parse_docx(path: str) -> List[Dict[str, Any]]:
    from docx import Document

    doc = Document(path)
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text)

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                full_text.append(row_text)

    combined = "\n".join(full_text)
    logger.info("Parsed DOCX %s: %d chars", path, len(combined))
    return [{"text": combined, "page": None}]


def parse_pptx(path: str) -> List[Dict[str, Any]]:
    from pptx import Presentation

    prs = Presentation(path)
    pages: List[Dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        texts.append(row_text)
        page_text = "\n".join(texts)
        if page_text:
            pages.append({"text": page_text, "page": i + 1})
    logger.info("Parsed PPTX %s: %d slides", path, len(pages))
    return pages